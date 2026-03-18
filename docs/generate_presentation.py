"""
OrgBlueprint Presentation Generator
Generates: OrgBlueprint_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A  slate-900
PANEL_BG     = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B  slate-800
BLUE_ACCENT  = RGBColor(0x3B, 0x82, 0xF6)   # #3B82F6  blue-500
BLUE_LIGHT   = RGBColor(0x60, 0xA5, 0xFA)   # #60A5FA  blue-400
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREY_TEXT    = RGBColor(0x94, 0xA3, 0xB8)   # #94A3B8  slate-400
GREEN_ACC    = RGBColor(0x34, 0xD3, 0x99)   # #34D399  emerald-400

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # truly blank

# ── Helper functions ──────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    fill_bg(slide)
    return slide


def fill_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name="Calibri", font_size=18, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_paragraph(tf, text, font_name="Calibri", font_size=16, bold=False, italic=False,
                  color=WHITE, align=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_bullet_textbox(slide, items, left, top, width, height,
                       font_size=18, title=None, title_size=22,
                       accent_bullets=False):
    """Add a textbox with optional title and bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = title
        run.font.name = "Calibri"
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = BLUE_LIGHT

    first = True
    for item in items:
        if title or not first:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        first = False

        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(5)

        if accent_bullets:
            # Bullet dot in accent color
            r1 = p.add_run()
            r1.text = "• "
            r1.font.name = "Calibri"
            r1.font.size = Pt(font_size)
            r1.font.color.rgb = BLUE_ACCENT
            r1.font.bold = True

            r2 = p.add_run()
            r2.text = item
            r2.font.name = "Calibri"
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = WHITE
        else:
            run = p.add_run()
            run.text = f"• {item}"
            run.font.name = "Calibri"
            run.font.size = Pt(font_size)
            run.font.color.rgb = WHITE

    return txBox


def slide_header(slide, number: str, label: str):
    """Tiny slide-number + section label in top-right corner."""
    add_textbox(slide, f"{number}  {label}",
                Inches(10.2), Inches(0.18), Inches(2.9), Inches(0.4),
                font_size=10, color=GREY_TEXT, align=PP_ALIGN.RIGHT)


def bottom_accent_bar(slide):
    """1-px blue line at very bottom."""
    add_rect(slide,
             Inches(0), H - Inches(0.08),
             W, Inches(0.08),
             fill_color=BLUE_ACCENT)


def left_accent_bar(slide, height=None):
    height = height or H
    add_rect(slide, Inches(0), Inches(0), Inches(0.1), height,
             fill_color=BLUE_ACCENT)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════════════════════
s1 = add_slide()

# Full-width decorative gradient strip at top
add_rect(s1, Inches(0), Inches(0), W, Inches(0.6), fill_color=BLUE_ACCENT)

# Large bottom accent bar
add_rect(s1, Inches(0), H - Inches(1.1), W, Inches(1.1), fill_color=PANEL_BG)

# Main title
tb = s1.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.5), Inches(1.8))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Practical Vibe Coding"
run.font.name = "Calibri"
run.font.size = Pt(60)
run.font.bold = True
run.font.color.rgb = WHITE

# Subtitle
add_textbox(s1, "Building a Production-Ready App with AI",
            Inches(0.8), Inches(2.85), Inches(9), Inches(0.7),
            font_size=28, color=BLUE_LIGHT, bold=False)

# Divider line
add_rect(s1, Inches(0.8), Inches(3.7), Inches(4.5), Inches(0.05),
         fill_color=BLUE_ACCENT)

# Tagline
add_textbox(s1, '"From a sales problem to a live product — in a weekend"',
            Inches(0.8), Inches(3.9), Inches(10), Inches(0.6),
            font_size=20, italic=True, color=GREY_TEXT)

# Presenter
add_textbox(s1, "Naveen",
            Inches(0.85), H - Inches(0.9), Inches(5), Inches(0.55),
            font_size=22, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — The Context: We Are Already in the Race
# ═══════════════════════════════════════════════════════════════════════════════
s2 = add_slide()
left_accent_bar(s2)
bottom_accent_bar(s2)
slide_header(s2, "02", "Context")

add_textbox(s2, "We Are Already in the Race",
            Inches(0.35), Inches(0.3), Inches(12), Inches(0.75),
            font_size=34, bold=True, color=WHITE)

add_rect(s2, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

# Hero statement panel
panel = add_rect(s2, Inches(0.35), Inches(1.35), Inches(12.6), Inches(1.1),
                 fill_color=PANEL_BG)

tb = s2.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(12.2), Inches(0.9))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = '"We are already in the middle of the AI era. This is not the warm-up. This is the race."'
run.font.name = "Calibri"
run.font.size = Pt(22)
run.font.bold = True
run.font.italic = True
run.font.color.rgb = BLUE_LIGHT

bullets = [
    "Speed matters. Curiosity matters. Adaptability matters.",
    "AI is not replacing us — it's making our experience, decisions, and output much faster and more productive.",
    "The question is not \"should I learn AI?\" — it's \"how fast can I start?\"",
]
add_bullet_textbox(s2, bullets,
                   Inches(0.35), Inches(2.7), Inches(12.6), Inches(3.5),
                   font_size=20, accent_bullets=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Where the Idea Came From
# ═══════════════════════════════════════════════════════════════════════════════
s3 = add_slide()
left_accent_bar(s3)
bottom_accent_bar(s3)
slide_header(s3, "03", "Origin Story")

add_textbox(s3, "Where the Idea Came From",
            Inches(0.35), Inches(0.3), Inches(12), Inches(0.75),
            font_size=34, bold=True, color=WHITE)

add_rect(s3, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

bullets = [
    "Working in sales, Naveen noticed a real gap",
    "Too many departments, too many handoffs, too many support calls — for something that could be answered with the right document and research",
    '"Why does it take so many people and so many iterations when the information already exists?"',
    "Thought: What if I could close that gap with something I build myself?",
    "Did not wait for a team, a budget, or a project plan — just started",
]
add_bullet_textbox(s3, bullets,
                   Inches(0.35), Inches(1.35), Inches(12.6), Inches(5.2),
                   font_size=20, accent_bullets=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — What is OrgBlueprint?
# ═══════════════════════════════════════════════════════════════════════════════
s4 = add_slide()
left_accent_bar(s4)
bottom_accent_bar(s4)
slide_header(s4, "04", "The Product")

add_textbox(s4, "What is OrgBlueprint?",
            Inches(0.35), Inches(0.3), Inches(12), Inches(0.75),
            font_size=34, bold=True, color=WHITE)

add_rect(s4, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

# Description panel
add_rect(s4, Inches(0.35), Inches(1.3), Inches(12.6), Inches(0.75),
         fill_color=PANEL_BG)
add_textbox(s4, "AI-powered Salesforce CRM blueprint generator — describe your business in plain English, get a full architecture plan.",
            Inches(0.55), Inches(1.38), Inches(12.2), Inches(0.65),
            font_size=17, italic=True, color=BLUE_LIGHT)

# Two columns
# Left
left_items = [
    "Product recommendations",
    "Architecture guidance",
    "Cost estimates",
    "Phased roadmap",
    "Document checklist",
]
add_bullet_textbox(s4, left_items,
                   Inches(0.35), Inches(2.2), Inches(5.5), Inches(3.0),
                   font_size=18, title="What it produces", accent_bullets=True)

# Right
right_items = [
    "Demo mode — instant, no API key required",
    "AI Enhanced — Orb asks 5 clarifying questions, LLM generates richer output",
]
add_bullet_textbox(s4, right_items,
                   Inches(6.2), Inches(2.2), Inches(6.8), Inches(1.5),
                   font_size=18, title="Two modes", accent_bullets=True)

audience_items = [
    "Founders, RevOps teams",
    "Salesforce admins & consultants",
    "Solution architects",
]
add_bullet_textbox(s4, audience_items,
                   Inches(6.2), Inches(4.05), Inches(6.8), Inches(2.0),
                   font_size=18, title="Built for", accent_bullets=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — The Tech Stack
# ═══════════════════════════════════════════════════════════════════════════════
s5 = add_slide()
left_accent_bar(s5)
bottom_accent_bar(s5)
slide_header(s5, "05", "Tech Stack")

add_textbox(s5, "The Tech Stack",
            Inches(0.35), Inches(0.3), Inches(12), Inches(0.75),
            font_size=34, bold=True, color=WHITE)

add_rect(s5, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

categories = [
    ("Frontend",  ["Next.js 14 (App Router) + React 18 + TypeScript",
                   "Tailwind CSS + shadcn/ui"]),
    ("Database & Auth", ["Prisma ORM + SQLite (local) / PostgreSQL (Vercel)",
                         "NextAuth v5"]),
    ("AI / LLM",  ["Anthropic Claude Sonnet → Gemini 2.0 Flash → Groq (fallback chain)",
                   "NVIDIA NIM (MiniMax M2.5) — chat widget fallback"]),
    ("Testing & Deploy", ["Playwright (E2E tests)",
                          "Vercel (deployment)"]),
]

col_w   = Inches(6.0)
col_gap = Inches(0.5)
row_h   = Inches(2.3)
start_y = Inches(1.35)

positions = [
    (Inches(0.35), start_y),
    (Inches(0.35) + col_w + col_gap, start_y),
    (Inches(0.35), start_y + row_h + Inches(0.15)),
    (Inches(0.35) + col_w + col_gap, start_y + row_h + Inches(0.15)),
]

for (lft, tp), (cat, items) in zip(positions, categories):
    add_rect(s5, lft, tp, col_w, row_h - Inches(0.1), fill_color=PANEL_BG)
    add_bullet_textbox(s5, items, lft + Inches(0.18), tp + Inches(0.12),
                       col_w - Inches(0.3), row_h - Inches(0.25),
                       font_size=16, title=cat, title_size=18, accent_bullets=True)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — What is Vibe Coding?
# ═══════════════════════════════════════════════════════════════════════════════
s6 = add_slide()
left_accent_bar(s6)
bottom_accent_bar(s6)
slide_header(s6, "06", "Vibe Coding")

add_textbox(s6, "What is Vibe Coding?",
            Inches(0.35), Inches(0.3), Inches(12), Inches(0.75),
            font_size=34, bold=True, color=WHITE)

add_rect(s6, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

bullets = [
    '"Vibe coding" = describing what you want in plain language, letting AI write it, reviewing and iterating',
    "Not about knowing every syntax or framework",
    "About having a clear idea + using AI as your co-developer",
    "The developer's role shifts: from writing every line  →  to directing, reviewing, and validating",
]
add_bullet_textbox(s6, bullets,
                   Inches(0.35), Inches(1.35), Inches(12.6), Inches(3.2),
                   font_size=19, accent_bullets=True)

# Pull-quote panel
add_rect(s6, Inches(0.35), Inches(4.75), Inches(12.6), Inches(1.85),
         fill_color=PANEL_BG)
# Left accent stripe on panel
add_rect(s6, Inches(0.35), Inches(4.75), Inches(0.12), Inches(1.85),
         fill_color=BLUE_ACCENT)

tb = s6.shapes.add_textbox(Inches(0.65), Inches(4.88), Inches(12.1), Inches(1.55))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = ('"The beauty of vibe coding is not that AI replaces thinking. '
            'The beauty is that it removes friction, speeds up experimentation, '
            'and allows ideas to become real much faster."')
run.font.name = "Calibri"
run.font.size = Pt(20)
run.font.bold = False
run.font.italic = True
run.font.color.rgb = BLUE_LIGHT


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — How It Was Built
# ═══════════════════════════════════════════════════════════════════════════════
s7 = add_slide()
left_accent_bar(s7)
bottom_accent_bar(s7)
slide_header(s7, "07", "The Process")

add_textbox(s7, "How It Was Built",
            Inches(0.35), Inches(0.3), Inches(12), Inches(0.75),
            font_size=34, bold=True, color=WHITE)

add_rect(s7, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

steps = [
    ("1", "Started with a rough idea — not a full spec"),
    ("2", "Used Claude Code as the AI pair programmer"),
    ("3", "Built in layers: core rules engine → Next.js app → auth → AI integration → tests → deployment"),
    ("4", "Every blocker was solved by prompting — not Googling for hours"),
    ("5", "Went from zero to production-deployed app in a weekend"),
    ("6", 'Key mindset: ship a working thing first, improve it after'),
]

y = Inches(1.4)
for num, text in steps:
    # Number circle (panel rect used as pseudo-circle)
    add_rect(s7, Inches(0.35), y, Inches(0.55), Inches(0.55), fill_color=BLUE_ACCENT)
    add_textbox(s7, num, Inches(0.35), y + Inches(0.03), Inches(0.55), Inches(0.5),
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s7, text, Inches(1.05), y, Inches(11.8), Inches(0.55),
                font_size=18, color=WHITE)
    y += Inches(0.82)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — Live Demo Walkthrough
# ═══════════════════════════════════════════════════════════════════════════════
s8 = add_slide()
left_accent_bar(s8)
bottom_accent_bar(s8)
slide_header(s8, "08", "Live Demo")

add_textbox(s8, "Live Demo Walkthrough",
            Inches(0.35), Inches(0.3), Inches(12), Inches(0.75),
            font_size=34, bold=True, color=WHITE)

add_rect(s8, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

steps_demo = [
    "Open OrgBlueprint at localhost:3000 (or Vercel URL)",
    "Type a business description in plain English",
    "Choose Demo mode → instant blueprint",
    "Or choose AI Enhanced → Orb asks clarifying questions → richer output",
    "Show all 6 dashboard tabs: Overview, Architecture, Data Model, Technical, Cost, Roadmap",
    "Show the floating AI chat widget",
    "Show PDF export",
]

y = Inches(1.4)
for i, step in enumerate(steps_demo, 1):
    add_rect(s8, Inches(0.35), y, Inches(0.55), Inches(0.55), fill_color=BLUE_ACCENT)
    add_textbox(s8, str(i), Inches(0.35), y + Inches(0.03), Inches(0.55), Inches(0.5),
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s8, step, Inches(1.05), y, Inches(11.8), Inches(0.55),
                font_size=18, color=WHITE)
    y += Inches(0.75)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — What OrgBlueprint Produces (table)
# ═══════════════════════════════════════════════════════════════════════════════
s9 = add_slide()
left_accent_bar(s9)
bottom_accent_bar(s9)
slide_header(s9, "09", "Output Sections")

add_textbox(s9, "What OrgBlueprint Produces",
            Inches(0.35), Inches(0.3), Inches(12), Inches(0.75),
            font_size=34, bold=True, color=WHITE)

add_rect(s9, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

rows = [
    ("Executive Snapshot",      "Confidence score, complexity, user count, primary focus"),
    ("Product Recommendations", "Recommended vs optional Salesforce products with rationale"),
    ("Architecture",            "OOTB vs custom, integration map, AppExchange picks"),
    ("Data Model",              "Entity cards, relationship diagram, automations"),
    ("Technical Blueprint",     "Platform-specific guidance per product"),
    ("Cost Estimate",           "Directional year-one range (not a quote)"),
    ("Roadmap",                 "Phased delivery + implementation checklist"),
    ("Document Checklist",      "BRD, data model, security model, test plan"),
]

# Header row
header_y = Inches(1.3)
add_rect(s9, Inches(0.35), header_y, Inches(4.0), Inches(0.42), fill_color=BLUE_ACCENT)
add_textbox(s9, "Section", Inches(0.45), header_y + Inches(0.04), Inches(3.8), Inches(0.36),
            font_size=15, bold=True, color=WHITE)
add_rect(s9, Inches(4.4), header_y, Inches(8.55), Inches(0.42), fill_color=BLUE_ACCENT)
add_textbox(s9, "What's inside", Inches(4.5), header_y + Inches(0.04), Inches(8.35), Inches(0.36),
            font_size=15, bold=True, color=WHITE)

row_h = Inches(0.65)
for i, (section, content) in enumerate(rows):
    ry = header_y + Inches(0.42) + i * row_h
    bg = PANEL_BG if i % 2 == 0 else DARK_BG
    add_rect(s9, Inches(0.35), ry, Inches(4.0), row_h, fill_color=bg)
    add_rect(s9, Inches(4.4), ry, Inches(8.55), row_h, fill_color=bg)

    add_textbox(s9, section, Inches(0.45), ry + Inches(0.08), Inches(3.8), row_h - Inches(0.1),
                font_size=14, bold=True, color=BLUE_LIGHT)
    add_textbox(s9, content, Inches(4.5), ry + Inches(0.08), Inches(8.35), row_h - Inches(0.1),
                font_size=14, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — What AI Actually Did in This Project
# ═══════════════════════════════════════════════════════════════════════════════
s10 = add_slide()
left_accent_bar(s10)
bottom_accent_bar(s10)
slide_header(s10, "10", "AI's Role")

add_textbox(s10, "What AI Actually Did in This Project",
            Inches(0.35), Inches(0.3), Inches(12), Inches(0.75),
            font_size=34, bold=True, color=WHITE)

add_rect(s10, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

ai_tasks = [
    "Wrote 80%+ of the boilerplate code",
    "Debugged TypeScript errors in seconds",
    "Suggested architecture patterns",
    "Generated the E2E test suite",
    "Fixed race conditions in React state",
    "Rewrote and improved the README",
    "Suggested better UX patterns for mobile",
]
add_bullet_textbox(s10, ai_tasks,
                   Inches(0.35), Inches(1.35), Inches(8.0), Inches(5.0),
                   font_size=19, accent_bullets=True)

# Pull-quote on the right
add_rect(s10, Inches(8.6), Inches(1.35), Inches(4.6), Inches(3.2),
         fill_color=PANEL_BG)
add_rect(s10, Inches(8.6), Inches(1.35), Inches(0.1), Inches(3.2),
         fill_color=BLUE_ACCENT)

tb = s10.shapes.add_textbox(Inches(8.85), Inches(1.6), Inches(4.1), Inches(2.8))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER

run = p.add_run()
run.text = '"I directed.\nAI executed.\nI reviewed.\nWe iterated."'
run.font.name = "Calibri"
run.font.size = Pt(24)
run.font.bold = True
run.font.italic = True
run.font.color.rgb = BLUE_LIGHT


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — Q&A: Technical Questions
# ═══════════════════════════════════════════════════════════════════════════════
s11 = add_slide()
left_accent_bar(s11)
bottom_accent_bar(s11)
slide_header(s11, "11", "Q&A — Technical")

add_textbox(s11, "Q&A — Technical Questions",
            Inches(0.35), Inches(0.3), Inches(12), Inches(0.75),
            font_size=34, bold=True, color=WHITE)

add_rect(s11, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

qa_tech = [
    ("How long did this actually take to build?",
     "Initial working prototype — about a weekend. Getting it production-quality with auth, tests, deployment, and AI fallbacks took a few more sessions."),
    ("What if Claude gives wrong code?",
     "You review it. Vibe coding doesn't mean you stop thinking — it means you spend less time typing and more time validating."),
    ("How do you handle the AI going off-track or hallucinations?",
     "Guardrails in the system prompt + deterministic fallback (rules engine). The app never returns bad data silently."),
    ("Is this production-ready or just a demo?",
     "It's live on Vercel, has a real database, auth, rate limiting, E2E tests, and a CI/CD pipeline. That's production by most definitions."),
    ("What was the hardest technical problem?",
     "TypeScript type mismatches between browser and Node.js environments. And getting the AI conversation flow deterministic so it never hangs."),
]

y = Inches(1.35)
row_h = Inches(1.0)
for q, a in qa_tech:
    add_rect(s11, Inches(0.35), y, Inches(12.6), row_h, fill_color=PANEL_BG)
    add_textbox(s11, f"Q: {q}", Inches(0.5), y + Inches(0.05), Inches(12.2), Inches(0.38),
                font_size=14, bold=True, color=BLUE_LIGHT)
    add_textbox(s11, f"A: {a}", Inches(0.5), y + Inches(0.42), Inches(12.2), Inches(0.52),
                font_size=13, color=GREY_TEXT)
    y += row_h + Inches(0.07)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — Q&A: Non-Technical Questions
# ═══════════════════════════════════════════════════════════════════════════════
s12 = add_slide()
left_accent_bar(s12)
bottom_accent_bar(s12)
slide_header(s12, "12", "Q&A — General")

add_textbox(s12, "Q&A — Non-Technical / General Questions",
            Inches(0.35), Inches(0.3), Inches(12.5), Inches(0.75),
            font_size=30, bold=True, color=WHITE)

add_rect(s12, Inches(0.35), Inches(1.15), Inches(12.7), Inches(0.05),
         fill_color=BLUE_ACCENT)

qa_gen = [
    ("How did you find the time?",
     "I use weekends — not because I have to, but because I'm genuinely curious. Even 2-3 hours on a Saturday compounds over months."),
    ("Isn't this just using ChatGPT to code?",
     "Claude Code is a dedicated AI coding agent — it reads your whole codebase, edits files, runs tests, and works like a pair programmer with full project context."),
    ("Should I be worried about AI replacing my job?",
     "AI is not replacing people — it's replacing tasks. The people who learn to use AI as a tool will move faster, build more, and create more value."),
    ("Do you need a CS degree to vibe code?",
     "No. You need a clear idea, patience to iterate, and enough curiosity to try. The AI handles the syntax."),
    ("Why Salesforce specifically?",
     "Expensive, complex ecosystem where the gap between 'what a business needs' and 'what an expert recommends' is huge — exactly where AI adds value."),
    ("How do you know the Salesforce recommendations are correct?",
     "Based on a deterministic rules engine with explicit guardrails — same logic a solutions architect would apply. AI enhances the narrative, not the product logic."),
]

y = Inches(1.35)
row_h = Inches(0.86)
for q, a in qa_gen:
    add_rect(s12, Inches(0.35), y, Inches(12.6), row_h, fill_color=PANEL_BG)
    add_textbox(s12, f"Q: {q}", Inches(0.5), y + Inches(0.04), Inches(12.2), Inches(0.35),
                font_size=13, bold=True, color=BLUE_LIGHT)
    add_textbox(s12, f"A: {a}", Inches(0.5), y + Inches(0.38), Inches(12.2), Inches(0.44),
                font_size=12, color=GREY_TEXT)
    y += row_h + Inches(0.05)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — The "Why Now" Message
# ═══════════════════════════════════════════════════════════════════════════════
s13 = add_slide()
bottom_accent_bar(s13)
slide_header(s13, "13", "Why Now")

# Full-bleed dark panel with a blue left border
add_rect(s13, Inches(0), Inches(0), Inches(0.18), H, fill_color=BLUE_ACCENT)

paragraphs = [
    ("We are already in the middle of the AI era.", True,  30, BLUE_LIGHT),
    ("This is not the warm-up. This is the race.", False, 26, WHITE),
    ("", False, 12, WHITE),
    ("Speed matters. Curiosity matters. Adaptability matters.", True, 22, WHITE),
    ("", False, 10, WHITE),
    ("You do not need to know every tool.", False, 20, GREY_TEXT),
    ("But you should be willing to try. To learn. To build.", False, 20, GREY_TEXT),
    ("", False, 10, WHITE),
    ("Open the tool. Test the idea. Build the prototype.", True, 22, BLUE_LIGHT),
    ("Because once you do that, AI stops being hype. It becomes leverage.", False, 20, WHITE),
    ("", False, 10, WHITE),
    ("So start small, but start now.", True, 28, GREEN_ACC),
]

tb = s13.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.6), Inches(6.9))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True

first = True
for text, bold, size, color in paragraphs:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(2)
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — The Beauty of Vibe Coding (Closing Quote)
# ═══════════════════════════════════════════════════════════════════════════════
s14 = add_slide()
bottom_accent_bar(s14)
slide_header(s14, "14", "Closing")

# Top strip
add_rect(s14, Inches(0), Inches(0), W, Inches(0.35), fill_color=BLUE_ACCENT)

# Large centered quote
tb = s14.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(4.2))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = ('"The beauty of vibe coding is not that AI replaces thinking.\n'
            'The beauty is that it removes friction, speeds up experimentation,\n'
            'and allows ideas to become real much faster."')
run.font.name = "Calibri"
run.font.size = Pt(30)
run.font.bold = False
run.font.italic = True
run.font.color.rgb = BLUE_LIGHT

# Divider
add_rect(s14, Inches(3.5), Inches(4.95), Inches(6.3), Inches(0.06),
         fill_color=BLUE_ACCENT)

# Credits
add_textbox(s14, "OrgBlueprint — built with Claude Code, Gemini, Next.js, and a lot of curiosity",
            Inches(0.8), Inches(5.15), Inches(11.7), Inches(0.5),
            font_size=16, color=GREY_TEXT, align=PP_ALIGN.CENTER)

add_textbox(s14, "github.com/iamnawin/orgblueprint-app",
            Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.45),
            font_size=16, bold=True, color=BLUE_ACCENT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — Thank You / Q&A
# ═══════════════════════════════════════════════════════════════════════════════
s15 = add_slide()
bottom_accent_bar(s15)

# Top strip
add_rect(s15, Inches(0), Inches(0), W, Inches(0.5), fill_color=BLUE_ACCENT)

# Large thank you
add_textbox(s15, "Thank You",
            Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.5),
            font_size=68, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Divider
add_rect(s15, Inches(3.5), Inches(2.85), Inches(6.3), Inches(0.07),
         fill_color=BLUE_ACCENT)

# Sub-text
add_textbox(s15,
            "Happy to take any questions\n— technical, non-technical, or just curious",
            Inches(0.8), Inches(3.05), Inches(11.7), Inches(1.1),
            font_size=26, italic=True, color=GREY_TEXT, align=PP_ALIGN.CENTER)

# Slide number label
slide_header(s15, "15", "Q&A")


# ═══════════════════════════════════════════════════════════════════════════════
#  Save
# ═══════════════════════════════════════════════════════════════════════════════
OUT = r"C:\Users\Naveen\OneDrive\Desktop\orgblueprint-clean\docs\OrgBlueprint_Presentation.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
